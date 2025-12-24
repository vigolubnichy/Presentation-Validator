from pptx import Presentation

def parse_presentation(path: str):

    prs = Presentation(path)
    slides_data = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        bullets = []
        fonts = set()
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_lines = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
            if text_lines:
                texts.extend(text_lines)
                if any(p.level > 0 or "•" in p.text for p in shape.text_frame.paragraphs):
                    bullets.append(text_lines)
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font and r.font.name:
                            fonts.add(r.font.name)

        slides_data.append({
            "slide_number": i,
            "texts": texts,
            "bullets": bullets,
            "fonts": list(fonts),
        })

    return slides_data
